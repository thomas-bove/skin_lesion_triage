"""Build presentation_1_topic_B_skin_lesion.pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util
from copy import deepcopy
import os

# ── paths ──────────────────────────────────────────────────────────────────
FIG = "/Users/thomasbove/Desktop/skin_lesion_project/skin_lesion/results/figures"
OUT = "/Users/thomasbove/Desktop/skin_lesion_project/skin_lesion/results/presentation/presentation_1_topic_B_skin_lesion.pptx"

# ── palette ────────────────────────────────────────────────────────────────
DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
MID    = RGBColor(0x16, 0x21, 0x3E)   # slide background
ACCENT = RGBColor(0x0F, 0x3C, 0x96)   # TU-blue
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF0, 0xF2, 0xF5)
MGREY  = RGBColor(0xCC, 0xD0, 0xD8)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x1A, 0x72, 0x48)

# ── slide size 16:9 ────────────────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # truly blank


# ══════════════════════════════════════════════════════════════════════════
# helpers
# ══════════════════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.RECTANGLE if False else 1,  # MSO_SHAPE.RECTANGLE = 1
        l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def txbox(slide, text, l, t, w, h,
          size=Pt(20), bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def section_header(slide, title):
    """Dark blue bar at top with white title."""
    add_rect(slide, 0, 0, W, Inches(0.85), fill=ACCENT)
    txbox(slide, title,
          Inches(0.35), Inches(0.08), W - Inches(0.7), Inches(0.7),
          size=Pt(32), bold=True, color=WHITE)


def footer(slide, num):
    add_rect(slide, 0, H - Inches(0.38), W, Inches(0.38), fill=DARK)
    txbox(slide,
          "Automatic Image Analysis — Presentation 1 — Topic B (HAM10000)",
          Inches(0.25), H - Inches(0.36), Inches(9), Inches(0.34),
          size=Pt(11), color=MGREY)
    txbox(slide, str(num),
          W - Inches(0.6), H - Inches(0.36), Inches(0.5), Inches(0.34),
          size=Pt(11), color=MGREY, align=PP_ALIGN.RIGHT)


def bg(slide):
    add_rect(slide, 0, 0, W, H, fill=WHITE)


def bullet_frame(slide, items, l, t, w, h,
                 size=Pt(20), color=DARK, bullet="•"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = size
        run.font.color.rgb = color
    return tb


def mini_table(slide, headers, rows, l, t, w, h,
               hdr_fill=ACCENT, hdr_color=WHITE,
               row_fills=(LGREY, WHITE),
               cell_size=Pt(16)):
    n_cols = len(headers)
    n_rows = len(rows)
    col_w = w // n_cols
    row_h = h // (n_rows + 1)

    # header row
    for ci, hdr in enumerate(headers):
        add_rect(slide, l + ci*col_w, t, col_w, row_h, fill=hdr_fill,
                 line=WHITE, line_w=Pt(1))
        txbox(slide, hdr, l + ci*col_w + Inches(0.04), t + Inches(0.02),
              col_w - Inches(0.08), row_h - Inches(0.04),
              size=cell_size, bold=True, color=hdr_color, align=PP_ALIGN.CENTER)

    # data rows
    for ri, row in enumerate(rows):
        fill = row_fills[ri % 2]
        for ci, cell in enumerate(row):
            add_rect(slide, l + ci*col_w, t + (ri+1)*row_h, col_w, row_h,
                     fill=fill, line=MGREY, line_w=Pt(0.5))
            txbox(slide, str(cell),
                  l + ci*col_w + Inches(0.04), t + (ri+1)*row_h + Inches(0.02),
                  col_w - Inches(0.08), row_h - Inches(0.04),
                  size=cell_size, color=DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
bg(s1)

# full-width accent bar
add_rect(s1, 0, Inches(1.8), W, Inches(3.4), fill=ACCENT)

# title
txbox(s1,
      "Skin Lesion Triage:\nA Transparent Bayesian Baseline on HAM10000",
      Inches(0.6), Inches(1.95), W - Inches(1.2), Inches(2.2),
      size=Pt(40), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# subtitle block
txbox(s1,
      "Thomas Bove  ·  Automatic Image Analysis  ·  TU Berlin  ·  Summer Semester 2026",
      Inches(0.6), Inches(4.15), W - Inches(1.2), Inches(0.55),
      size=Pt(22), color=WHITE, align=PP_ALIGN.CENTER)

# research question box
add_rect(s1, Inches(1.0), Inches(5.1), W - Inches(2.0), Inches(1.5),
         fill=LGREY, line=ACCENT, line_w=Pt(2))
txbox(s1,
      "Research question: Can a class-conditional GMM with an asymmetric cost matrix "
      "produce clinically acceptable triage decisions on HAM10000, "
      "and is its confidence well-calibrated?",
      Inches(1.15), Inches(5.2), W - Inches(2.3), Inches(1.3),
      size=Pt(20), color=DARK, align=PP_ALIGN.CENTER)

# top stripe
add_rect(s1, 0, 0, W, Inches(0.25), fill=DARK)
add_rect(s1, 0, H - Inches(0.25), W, Inches(0.25), fill=DARK)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Dataset, Features, Baseline Pipeline
# ══════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
bg(s2)
section_header(s2, "Dataset, Features & Classical Baseline")
footer(s2, 2)

# ── left column: dataset facts ────────────────────────────────────────────
txbox(s2, "Dataset & Splits", Inches(0.35), Inches(1.05), Inches(5.5), Inches(0.45),
      size=Pt(22), bold=True, color=ACCENT)

mini_table(s2,
           ["Split", "Total", "Melanoma", "Benign"],
           [["Train", "1 335", "668  (P=0.500)", "667"],
            ["Val",   "445",   "222  (P=0.499)", "223"],
            ["Test",  "446",   "223",             "223"]],
           Inches(0.35), Inches(1.5), Inches(5.9), Inches(1.9),
           cell_size=Pt(17))

txbox(s2,
      "Binary task: melanoma (mel) vs. benign nevus (nv)\n"
      "Balanced subset: 1 113 mel + 1 113 nv  →  2 226 images total\n"
      "Stratified 60 / 20 / 20 split  ·  random seed = 42",
      Inches(0.35), Inches(3.45), Inches(5.9), Inches(0.9),
      size=Pt(17), color=DARK)

# ── right column: features + pipeline ────────────────────────────────────
txbox(s2, "Features & Pipeline", Inches(6.6), Inches(1.05), Inches(6.4), Inches(0.45),
      size=Pt(22), bold=True, color=ACCENT)

bullet_frame(s2,
             ["Feature: HSV per-channel histograms (H 32 + S 32 + V 32 = 96-dim), L1-normalised",
              "Class-conditional GMM fit by EM (full covariance)",
              "K per class selected by BIC on training set (K=1…6)",
              "Posterior via Bayes rule (log-space, numerically stable)",
              "Decision threshold derived from asymmetric cost matrix"],
             Inches(6.6), Inches(1.55), Inches(6.3), Inches(2.6),
             size=Pt(18), color=DARK)

add_rect(s2, Inches(6.6), Inches(4.25), Inches(6.3), Inches(0.9),
         fill=LGREY, line=MGREY, line_w=Pt(1))
txbox(s2,
      "Modern comparison plan: deep CNN (EfficientNet-B0) "
      "with temperature scaling — Presentation 2",
      Inches(6.7), Inches(4.3), Inches(6.1), Inches(0.8),
      size=Pt(17), color=DARK)

add_rect(s2, Inches(6.25), Inches(1.0), Pt(1.5), Inches(5.5), fill=MGREY)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Cost Matrix & Threshold Experiment
# ══════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
bg(s3)
section_header(s3, "Cost Matrix & Threshold Experiment")
footer(s3, 3)

# ── cost matrix ──────────────────────────────────────────────────────────
txbox(s3, "Loss matrix  λ(action | true class)", Inches(0.35), Inches(1.0),
      Inches(5.5), Inches(0.45), size=Pt(20), bold=True, color=ACCENT)

# hand-built cost matrix with coloured cells
cell_w = Inches(1.55)
cell_h = Inches(0.65)
ox, oy = Inches(0.35), Inches(1.5)

# header cells
for ci, txt in enumerate(["", "True: Melanoma", "True: Benign"]):
    fill = ACCENT if ci > 0 else MID
    add_rect(s3, ox + ci*cell_w, oy, cell_w, cell_h, fill=fill,
             line=WHITE, line_w=Pt(1))
    txbox(s3, txt, ox + ci*cell_w + Inches(0.05), oy + Inches(0.05),
          cell_w - Inches(0.1), cell_h - Inches(0.1),
          size=Pt(16), bold=(ci>0), color=WHITE, align=PP_ALIGN.CENTER)

data_rows = [
    ("Predict: Melanoma", "0  (correct)", "1  (FP)"),
    ("Predict: Benign",   "5  (FN !!)",   "0  (correct)"),
]
row_fills_cm = [LGREY, RGBColor(0xFF, 0xEB, 0xEB)]
for ri, (label, c1, c2) in enumerate(data_rows):
    add_rect(s3, ox, oy + (ri+1)*cell_h, cell_w, cell_h, fill=ACCENT,
             line=WHITE, line_w=Pt(1))
    txbox(s3, label, ox + Inches(0.05), oy + (ri+1)*cell_h + Inches(0.05),
          cell_w - Inches(0.1), cell_h - Inches(0.1),
          size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for ci, val in enumerate([c1, c2]):
        col = RED if "FN" in val else (DARK if "correct" in val else DARK)
        add_rect(s3, ox + (ci+1)*cell_w, oy + (ri+1)*cell_h, cell_w, cell_h,
                 fill=row_fills_cm[ri], line=MGREY, line_w=Pt(0.75))
        txbox(s3, val,
              ox + (ci+1)*cell_w + Inches(0.05),
              oy + (ri+1)*cell_h + Inches(0.05),
              cell_w - Inches(0.1), cell_h - Inches(0.1),
              size=Pt(17), bold=("FN" in val), color=col,
              align=PP_ALIGN.CENTER)

txbox(s3,
      "λ(FN) = 5 · λ(FP)  →  threshold:  P(ω_M | x) > 1/6 ≈ 0.167   (vs MAP: 0.5)",
      Inches(0.35), Inches(3.05), Inches(5.6), Inches(0.55),
      size=Pt(18), bold=True, color=ACCENT)

# ── threshold experiment ──────────────────────────────────────────────────
txbox(s3, "Threshold experiment — test set (n = 446)",
      Inches(6.3), Inches(1.0), Inches(6.7), Inches(0.45),
      size=Pt(20), bold=True, color=ACCENT)

mini_table(s3,
           ["Rule", "Threshold", "TP", "FP", "TN", "FN", "Recall", "FPR", "Cost"],
           [["MAP",            "0.500", "187", "56", "167", "36", "0.839", "0.251", "236"],
            ["Cost-sensitive", "1/6",   "188", "59", "164", "35", "0.843", "0.265", "234"]],
           Inches(6.3), Inches(1.5), Inches(6.65), Inches(1.75),
           cell_size=Pt(15))

txbox(s3,
      "Cost-sensitive rule: −1 FN, +3 FP vs MAP.\n"
      "Small gain because GMM posteriors are strongly bimodal — "
      "very few samples fall in the swing band [1/6, 0.5).\n"
      "Direction matches theory (slides 9–10); magnitude limited by feature representation.",
      Inches(6.3), Inches(3.3), Inches(6.65), Inches(1.4),
      size=Pt(17), color=DARK)

# embed confusion matrix figure
s3.shapes.add_picture(
    os.path.join(FIG, "confusion_matrices.png"),
    Inches(6.3), Inches(4.75), Inches(6.65), Inches(2.35))


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — BIC Model Selection & Calibration
# ══════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
bg(s4)
section_header(s4, "BIC Model Selection & Calibration")
footer(s4, 4)

# ── left: BIC ─────────────────────────────────────────────────────────────
txbox(s4, "BIC — K selection (training set)",
      Inches(0.35), Inches(1.0), Inches(5.8), Inches(0.4),
      size=Pt(20), bold=True, color=ACCENT)

s4.shapes.add_picture(
    os.path.join(FIG, "bic_curves.png"),
    Inches(0.35), Inches(1.45), Inches(5.8), Inches(2.55))

mini_table(s4,
           ["Class", "K chosen", "Covariance", "Converged"],
           [["Melanoma", "1", "full", "Yes (2 iter)"],
            ["Benign",   "2", "full", "Yes (13 iter)"]],
           Inches(0.35), Inches(4.1), Inches(5.8), Inches(1.2),
           cell_size=Pt(15))

txbox(s4,
      "K=1 for melanoma: BIC penalty (~31 k per component in 96D) "
      "exceeds LL gain at n=668.\nStable across 5 seeds (42, 7, 123, 99, 314).",
      Inches(0.35), Inches(5.35), Inches(5.8), Inches(0.75),
      size=Pt(16), color=DARK)

# ── right: calibration ────────────────────────────────────────────────────
txbox(s4, "Calibration — reliability diagram (test set, 15 bins)",
      Inches(6.45), Inches(1.0), Inches(6.55), Inches(0.4),
      size=Pt(20), bold=True, color=ACCENT)

s4.shapes.add_picture(
    os.path.join(FIG, "reliability_diagram_calibrated.png"),
    Inches(6.45), Inches(1.45), Inches(6.55), Inches(3.05))

mini_table(s4,
           ["", "ECE", "T"],
           [["Before scaling", "0.206", "1.0 (none)"],
            ["After scaling",  "0.129", "35.7 (val-fit)"]],
           Inches(6.45), Inches(4.6), Inches(6.55), Inches(1.1),
           cell_size=Pt(15))

txbox(s4,
      "Bimodal posteriors: bars only near 0 and 1, middle bins empty.\n"
      "Model is overconfident. Temperature scaling (T≈36) gives partial improvement;\n"
      "high ECE persists — 96D HSV histogram may not support calibrated GMM.",
      Inches(6.45), Inches(5.78), Inches(6.55), Inches(0.85),
      size=Pt(15), color=DARK)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Evaluation Summary, Blockers, Next Steps
# ══════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
bg(s5)
section_header(s5, "Evaluation Summary, Blockers & Next Steps")
footer(s5, 5)

# ── top: headline metrics ─────────────────────────────────────────────────
txbox(s5, "Headline metrics — test set",
      Inches(0.35), Inches(1.0), Inches(7.0), Inches(0.4),
      size=Pt(20), bold=True, color=ACCENT)

mini_table(s5,
           ["Metric", "MAP (θ=0.5)", "Cost-sens (θ=1/6)"],
           [["ROC-AUC",          "0.835",  "0.835"],
            ["Sensitivity",      "0.839",  "0.843"],
            ["Specificity",      "0.749",  "0.735"],
            ["FN (missed mel.)", "36",     "35"],
            ["ECE (raw / T-sc)", "0.206",  "0.129"]],
           Inches(0.35), Inches(1.45), Inches(7.0), Inches(3.0),
           cell_size=Pt(16))

# ── ROC curve ────────────────────────────────────────────────────────────
s5.shapes.add_picture(
    os.path.join(FIG, "roc_curve.png"),
    Inches(7.6), Inches(0.9), Inches(5.4), Inches(3.6))

# ── blockers ─────────────────────────────────────────────────────────────
txbox(s5, "Current blockers",
      Inches(0.35), Inches(4.6), Inches(6.5), Inches(0.4),
      size=Pt(20), bold=True, color=RED)

bullet_frame(s5,
             ["96-D HSV histogram: BIC selects K=1 for melanoma "
              "(penalty ~31 k/component dominates at n=668)",
              "Bimodal posteriors → ECE=0.21 raw; temperature scaling reduces to 0.13 "
              "but residual miscalibration remains",
              "Cost-sensitive threshold gains only −1 FN (swing band nearly empty)",
              "Feature representation may be too weak to discriminate mel vs nv"],
             Inches(0.35), Inches(5.05), Inches(7.2), Inches(2.0),
             size=Pt(17), color=DARK)

# ── next steps ────────────────────────────────────────────────────────────
add_rect(s5, Inches(7.6), Inches(4.55), Inches(5.4), Inches(2.55),
         fill=LGREY, line=ACCENT, line_w=Pt(1.5))
txbox(s5, "Presentation 2 will add:",
      Inches(7.75), Inches(4.65), Inches(5.1), Inches(0.4),
      size=Pt(18), bold=True, color=ACCENT)
bullet_frame(s5,
             ["EfficientNet-B0 deep features (1280-dim)",
              "Trainable classification head",
              "Temperature scaling comparison (GMM vs CNN)",
              "Side-by-side ROC, ECE, FN/FP table"],
             Inches(7.75), Inches(5.1), Inches(5.1), Inches(1.85),
             size=Pt(17), color=DARK)


# ── save ──────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
